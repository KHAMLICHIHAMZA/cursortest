from __future__ import annotations

from datetime import date
from pathlib import Path

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches


ROOT = Path("C:/Projects/MALOC")
EVIDENCE = ROOT / "docs" / "evidence" / "live-role-screens"
LOGO = ROOT / "mobile-agent" / "assets" / "icon.png"

OUTPUT_PDF = ROOT / "docs" / "MALOC_PRESENTATION_CLIENT_COMPLETE_2026-03-07.pdf"
OUTPUT_PPT = ROOT / "docs" / "MALOC_PRESENTATION_CLIENT_COMPLETE_2026-03-07.pptx"

SECTIONS = [
    ("public", "Version Publique"),
    ("super_admin", "Version Admin (Super Admin)"),
    ("company_admin", "Version Company (Direction Entreprise)"),
    ("agency_manager", "Version Agence (Manager)"),
    ("agent", "Version Agent (Terrain)"),
]


def label(stem: str) -> str:
    return stem.replace("_", " ").replace("-", " ").strip().title()


def collect_images(section_key: str) -> list[Path]:
    folder = EVIDENCE / section_key
    if not folder.exists():
        return []
    return sorted([p for p in folder.glob("*.png") if p.is_file()])


def build_pdf() -> None:
    pdf = FPDF(orientation="L", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=10)

    # Cover
    pdf.add_page()
    if LOGO.exists():
        pdf.image(str(LOGO), x=10, y=8, w=20)
    pdf.set_font("Helvetica", "B", 28)
    pdf.ln(20)
    pdf.cell(0, 12, "MALOC - Presentation Client Complete", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 14)
    pdf.cell(0, 10, f"Date: {date.today().isoformat()}", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)
    pdf.multi_cell(
        0,
        9,
        "Document visuel complet a destination client. "
        "Cette version regroupe toutes les captures disponibles par version utilisateur.",
    )

    total = 0
    for key, _ in SECTIONS:
        total += len(collect_images(key))
    pdf.ln(2)
    pdf.set_font("Helvetica", "B", 13)
    pdf.cell(0, 8, f"Total captures integrees: {total}", new_x="LMARGIN", new_y="NEXT")

    for key, section_title in SECTIONS:
        images = collect_images(key)
        # Section page
        pdf.add_page()
        if LOGO.exists():
            pdf.image(str(LOGO), x=10, y=8, w=14)
        pdf.set_font("Helvetica", "B", 22)
        pdf.ln(14)
        pdf.cell(0, 12, section_title, new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 13)
        pdf.cell(0, 8, f"Nombre d'ecrans: {len(images)}", new_x="LMARGIN", new_y="NEXT")

        for image in images:
            pdf.add_page()
            if LOGO.exists():
                pdf.image(str(LOGO), x=10, y=8, w=10)
            pdf.set_font("Helvetica", "B", 14)
            pdf.ln(6)
            pdf.cell(0, 8, label(image.stem), new_x="LMARGIN", new_y="NEXT")
            # Fit screenshot in landscape page with margins
            pdf.image(str(image), x=10, y=24, w=277, h=180, keep_aspect_ratio=True)
            pdf.set_y(206)
            pdf.set_font("Helvetica", "", 10)
            pdf.cell(0, 6, f"Fichier: {image.name}")

    pdf.output(str(OUTPUT_PDF))


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_section_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_image_slide(prs: Presentation, image_path: Path, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    # place image to maximize viewport below title
    left = Inches(0.4)
    top = Inches(1.0)
    width = Inches(12.5)
    height = Inches(6.0)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def build_ppt() -> None:
    prs = Presentation()
    add_title_slide(
        prs,
        "MALOC - Presentation Client Complete",
        f"Toutes versions | Captures live | {date.today().isoformat()}",
    )

    for key, section_title in SECTIONS:
        images = collect_images(key)
        add_section_slide(prs, section_title, f"{len(images)} ecrans presentes")
        for image in images:
            add_image_slide(prs, image, label(image.stem))

    prs.save(str(OUTPUT_PPT))


def main() -> None:
    build_pdf()
    build_ppt()
    print(f"Generated PDF: {OUTPUT_PDF}")
    print(f"Generated PPT: {OUTPUT_PPT}")


if __name__ == "__main__":
    main()
