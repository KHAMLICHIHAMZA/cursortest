from pptx import Presentation


def add_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0


def main() -> None:
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "MALOC - Synthese de livraison"
    title_slide.placeholders[1].text = "Charges, booking, permissions, mobile, UAT\nMise a jour: 2026-03-07"

    add_bullets(
        prs,
        "Objectifs de la livraison",
        [
            "Fiabiliser les cycles booking et post-checkout",
            "Renforcer controle des droits par role",
            "Professionnaliser Charges & Depenses et KPI",
            "Stabiliser flux mobile check-in/check-out",
        ],
    )

    add_bullets(
        prs,
        "Charges & Depenses - apports",
        [
            "Portee VEHICLE/AGENCY/COMPANY + centre de cout",
            "Categories explicites non-vehicule (salary, office rent...)",
            "Filtrage intelligent des categories selon contexte",
            "Export CSV robuste compatible Excel (UTF-8/CRLF)",
        ],
    )

    add_bullets(
        prs,
        "Booking et gouvernance des droits",
        [
            "AGENT bloque sur create/update/delete booking",
            "AGENT conserve check-in/check-out",
            "Cloture financiere reservee aux roles agence autorises",
            "Notification in-app de cloture en attente apres checkout",
        ],
    )

    add_bullets(
        prs,
        "Qualite operationnelle",
        [
            "Deduplication notifications (pas de doublons)",
            "Activation email et profil obligatoire pour roles agence",
            "KPI enrichi: charges par centre de cout et allocation partagee",
            "Scripts de simulation fonctionnelle et E2E role AGENT",
        ],
    )

    add_bullets(
        prs,
        "Mobile agent",
        [
            "Correction WebView natif pour signature",
            "Hydratation donnees check-in dans ecran check-out",
            "Suppression du cas odometerStart non defini",
            "Flux terrain plus fiable pour agents",
        ],
    )

    add_bullets(
        prs,
        "Plan de preuves et UAT",
        [
            "Checklist UAT dediee (roles, booking, charges, mobile)",
            "Plan de captures ecran structure par module",
            "Preuve CSV et preuves API 403/succes attendues",
            "Gate release: 0 blocant, ecarts mineurs documentes",
        ],
    )

    add_bullets(
        prs,
        "Prochaines etapes",
        [
            "Executer UAT en environnement cible",
            "Inserer captures finales dans le dossier preuves",
            "Valider resultat avec metier et operations agence",
            "Programmer fenetre de mise en production",
        ],
    )

    prs.save("docs/PRESENTATION_LIVRAISON_2026-03-07.pptx")


if __name__ == "__main__":
    main()
